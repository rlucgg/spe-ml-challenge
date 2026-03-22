"""Generate best-in-class presentation for SPE GCS 2026 ML Challenge."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────────
NAVY = RGBColor(0x0D, 0x2B, 0x4E)
TEAL = RGBColor(0x00, 0x7B, 0x83)
GOLD = RGBColor(0xD4, 0xA0, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_BG = RGBColor(0xE8, 0xF4, 0xF8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _accent_bar(slide, y=Inches(0), height=Inches(0.06), color=TEAL):
    """Add a thin coloured accent bar across the full width."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), y, SLIDE_W, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _side_stripe(slide, color=NAVY, width=Inches(0.35)):
    """Add a left-edge vertical stripe."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), width, SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(slide, left, top, width, height, text,
              size=16, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
              font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def _add_bullets(slide, left, top, width, height, items,
                 size=15, color=DARK_TEXT, indent_size=13,
                 indent_color=MED_GRAY, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_sub = item.startswith("  ")
        p.text = item.strip()
        p.font.size = Pt(indent_size if is_sub else size)
        p.font.color.rgb = indent_color if is_sub else color
        p.font.name = "Calibri"
        p.level = 1 if is_sub else 0
        p.space_after = spacing


def _stat_box(slide, left, top, number, label, color=TEAL):
    """Add a bold metric box: big number on top, label below."""
    w, h = Inches(2.4), Inches(1.3)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_BG
    bg.line.fill.background()
    _add_text(slide, left, top + Inches(0.1), w, Inches(0.7),
              number, size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    _add_text(slide, left, top + Inches(0.7), w, Inches(0.5),
              label, size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ── Slide builders ──────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    _accent_bar(slide, y=Inches(3.3), height=Inches(0.05), color=GOLD)
    _add_text(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(1.5),
              "Agentic AI for Drilling\nOperational Intelligence",
              size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1.2), Inches(3.6), Inches(11), Inches(1.0),
              "SPE GCS 2026 ML Challenge\nAnalyzing the Equinor Volve Field Dataset",
              size=22, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)
    _add_text(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(0.6),
              "Rong Lu  |  March 2026", size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)


def slide_challenge(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "The Challenge", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              "Build an evidence-based AI agent for drilling operational intelligence",
              size=14, color=MED_GRAY)
    _add_bullets(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0), [
        "Analyze the Equinor Volve Field drilling dataset",
        "Answer 6 categories of operational questions:",
        "  1. Drilling Phase Identification",
        "  2. Time & Efficiency Analysis (NPT)",
        "  3. Section & ROP Performance",
        "  4. BHA Configuration Effectiveness",
        "  5. Operational Issues & Root Causes",
        "  6. Synthesis & Recommendations",
    ])
    # Right column: key evaluation criteria
    _add_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
              "Evaluation Criteria", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5), [
        "Quality of REASONING (most important)",
        "Evidence from BOTH data AND reports",
        "Consistency across answers",
        "Clear assumptions & uncertainty handling",
        "Transparency & reproducibility",
        '"Complexity alone will not be rewarded"',
    ], size=14, color=MED_GRAY)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "System Architecture", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              "Pure OpenAI SDK + DuckDB + ChromaDB  |  No LangChain, no framework bloat",
              size=14, color=MED_GRAY)
    # Flow description
    _add_bullets(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5), [
        "User Question",
        "  --> GPT-5.4 mini (tool calling, max 10 rounds)",
        "  --> 12 Agent Tools (see right)",
        "  --> Structured Answer with Evidence",
        "",
        "Key features:",
        "  Evidence Trace mode (--trace) shows every",
        "    tool call, args, timing, and data sources",
        "  Retry with exponential backoff for API resilience",
        "  Output validation: checks all 6 required sections",
    ], size=14)

    # Right column: tool list
    _add_text(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5),
              "12 Agent Tools", size=18, bold=True, color=TEAL)
    tools = [
        ("query_drilling_data", "SQL on 12 DuckDB tables"),
        ("search_daily_reports", "Semantic search (ChromaDB)"),
        ("get_well_overview", "Well metadata & formations"),
        ("get_drilling_phases", "Hole-size + activity-code phases"),
        ("compute_efficiency", "NPT sub-classification, ROP"),
        ("compare_wells", "Side-by-side + production data"),
        ("get_bha_configurations", "WITSML BHA + mudlog params"),
        ("identify_issues", "Root cause + mud correlation"),
        ("get_formation_context", "Geological depth context"),
        ("get_field_benchmarks", "Cross-well rankings (5 modes)"),
        ("get_ddr_narrative", "Guaranteed DDR text by date"),
        ("generate_depth_time_plot", "Depth-vs-time visualization"),
    ]
    tf = _add_text(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(0.3),
                   "", size=11)
    for name, desc in tools:
        p = tf.add_paragraph()
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(1)
        run1 = p.add_run()
        run1.text = f"{name:<28}"
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.name = "Consolas"
        run2 = p.add_run()
        run2.text = f" {desc}"
        run2.font.size = Pt(10)
        run2.font.color.rgb = MED_GRAY


def slide_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "Data Integration", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              "12 DuckDB tables from 4 data sources  |  Ingestion: ~3 minutes",
              size=14, color=MED_GRAY)

    # Stat boxes row
    _stat_box(slide, Inches(0.8),  Inches(1.6), "1,759",  "DDR XML Files\n(zero parse errors)")
    _stat_box(slide, Inches(3.5),  Inches(1.6), "23,447", "Activity Records")
    _stat_box(slide, Inches(6.2),  Inches(1.6), "2,882",  "MudLog Intervals\n(ROP/WOB/RPM/lithology)")
    _stat_box(slide, Inches(8.9),  Inches(1.6), "11,134", "WITSML Messages")
    _stat_box(slide, Inches(11.6), Inches(1.6), "~38K",   "ChromaDB Docs\n(DDR + WITSML)")

    # Data details
    _add_bullets(slide, Inches(0.8), Inches(3.3), Inches(5.8), Inches(4.0), [
        "DDR Daily Drilling Reports (core NLP source)",
        "  23,447 activities with timestamps & codes",
        "  2,271 fluid records (mud type, density, PV, YP)",
        "  1,726 directional survey stations",
        "  26,965 text documents (summaries + comments)",
        "",
        "WITSML Real-Time Data (measured drilling parameters)",
        "  2,882 mudlog intervals: actual ROP, WOB, torque, RPM",
        "  161 official BHA runs across 14 wells",
        "  4,217 trajectory stations (MD, TVD, incl, azi, DLS)",
    ], size=13)

    _add_bullets(slide, Inches(7.0), Inches(3.3), Inches(5.8), Inches(4.0), [
        "Production Data",
        "  15,634 daily records: oil, gas, water, pressure",
        "  7 wells: F-1C, F-11, F-12, F-14, F-15D, F-4, F-5",
        "",
        "Geological Data",
        "  409 formation tops across 26 wells",
        "  48 perforation interval records",
        "",
        "Data Quality",
        "  Sentinel values (-999.xx) filtered at parse time",
        "  Outlier filtering: RPM<=300, WOB<=500kN, ROP<=200 m/hr",
        "  Unit conversions: m/s->m/hr, N->kN, rad->deg, kg/m3->sg",
    ], size=13)


def slide_evidence_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "Evidence-First Design", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              "Every answer cites BOTH structured data AND daily report quotes",
              size=14, color=MED_GRAY)

    _add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
              "Mandatory Answer Structure (6 Sections)", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), [
        "1. Answer  (clear, concise conclusion)",
        "2. Evidence from Drilling Data  (specific measurements)",
        "3. Evidence from Daily Reports  (DDR quotes with dates)",
        "4. Reasoning  (numbered step-by-step chain)",
        "5. Assumptions  (explicitly stated)",
        "6. Confidence & Uncertainty  (HIGH / MEDIUM / LOW)",
    ], size=15)

    _add_text(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.5),
              "Key Differentiators", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5), [
        "WITSML mudlog = actual measured ROP/WOB/RPM",
        "  (not estimated from daily depth progress)",
        "",
        "Hole-size transitions define phases",
        "  (not just generic activity code classification)",
        "",
        "Formation context links ROP to geology",
        "  (why is this section slow? check the lithology)",
        "",
        "Statistical mud-problem correlation",
        "  (problem-day vs normal-day property comparison)",
        "",
        "NPT sub-classified by comment analysis",
        "  (11 categories instead of generic 'Other NPT')",
    ], size=13)


def slide_example(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "Example: Phase Identification", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              'Q: "Identify the major drilling phases for well 15/9-F-11 T2"',
              size=14, color=MED_GRAY)

    _add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
              "Agent Tool Calls", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.0), [
        "1. get_drilling_phases('15_9_F_11_T2')",
        "2. get_ddr_narrative(well, date_from, date_to)",
        "3. get_formation_context('15_9_F_11_T2')",
    ], size=14)

    _add_text(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(0.4),
              "Answer (3 major hole sections identified)", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0), [
        'Phase 1: Surface 26" (Mar 24 - Apr 13)',
        "  306m - 1,365m MD | avg ROP 41.6 m/hr",
        'Phase 2: Intermediate 17.5" (Apr 14 - Apr 28)',
        "  1,400m - 2,577m MD | avg ROP 29.2 m/hr",
        'Phase 3: Reservoir 8.5" (Apr 29 - May 9)',
        "  2,907m - 4,562m MD | avg ROP 19.0 m/hr",
    ], size=13)

    # Right column: evidence excerpts
    _add_text(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.4),
              "Evidence from Daily Reports", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.5), [
        'DDR 2013-04-14: "Drilled 17 1/2" hole from',
        '  1368m to 1400m with 3675 l/min, 188 bar,',
        '  60-140 RPM, 7-15 kNm, WOB 5 Ton"',
        "",
        'DDR 2013-04-22: "Drilled 17 1/2" hole to',
        '  section TD at 2574m. POOH and L/O',
        '  17 1/2" steerable BHA."',
        "",
        "Confidence: HIGH",
        "  Hole size changes confirmed by activity",
        "  codes, depth progression, and DDR summaries",
    ], size=13, color=MED_GRAY)


def slide_field_benchmarks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "Cross-Well Analysis", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              "Field-wide benchmarking across 26 wellbore sections",
              size=14, color=MED_GRAY)

    _add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
              "get_field_benchmarks (5 Modes)", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), [
        "daily_progress",
        "  Ranks well-sections by avg drilled m/day",
        "section_performance",
        "  Difficulty index from WOB, torque, ROP z-scores",
        "gas_response",
        "  Peak methane/ethane in target formation window",
        "risk",
        "  Composite score: stuck/loss/kick mentions + severity",
        "production_summary",
        "  Cumulative oil/gas/water with normalized well names",
    ], size=14)

    _add_text(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.4),
              "Example Results", size=18, bold=True, color=TEAL)
    _add_bullets(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5), [
        "Hardest section to drill (by difficulty index):",
        '  F-11 B 12.25" section (difficulty 3.07)',
        "  avg ROP 25.0 m/hr | WOB 112.1 kN",
        "",
        "Highest gas response in Hugin Fm:",
        "  F-1 C: max CH4 35,013 ppm at 3229-4004m MD",
        "",
        "Highest risk well:",
        "  F-15: risk score 168.5 (51 severe mentions)",
        "",
        "Top oil producer:",
        "  F-12: cum oil 4.58M Sm3 over 3,056 prod days",
    ], size=13, color=MED_GRAY)


def slide_trace(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "Evidence Trace: Transparent Reasoning", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              "python -m src.main ask '...' --trace   |   Turns the black box into a glass box",
              size=14, color=MED_GRAY)

    _add_bullets(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.0), [
        "## Evidence Trace",
        "",
        "### Step 1: get_bha_configurations(well='15_9_F_11_T2')",
        "  Retrieved: 3,475 chars | Duration: 0.26s",
        "  Summary: 12 official BHA runs, ROP by section, lithology...",
        "",
        "### Step 2: get_ddr_narrative(well='15_9_F_11_T2', date_from='2013-04-11', date_to='2013-04-26')",
        "  Retrieved: 7,134 chars | Duration: 0.02s",
        '  Summary: 16 DDR reports, 67 activities; "Drilled 17 1/2" hole to section TD..."',
        "",
        "Total tool calls: 3 | Total evidence retrieval time: 0.29s",
        "Data sources used: ddr_activities, ddr_status, witsml_bha_runs, witsml_mudlog",
    ], size=14)

    _add_text(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
              "Judges can see EXACTLY how the agent reasoned. Maximum transparency.",
              size=16, bold=True, color=TEAL)


def slide_design_decisions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _side_stripe(slide)
    _accent_bar(slide, y=Inches(7.44))
    _add_text(slide, Inches(0.8), Inches(0.35), Inches(11), Inches(0.7),
              "Key Design Decisions", size=30, bold=True, color=NAVY)
    _add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
              '"Complexity alone will not be rewarded" -- Judges',
              size=14, color=MED_GRAY)

    _add_bullets(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.5), [
        "OpenAI SDK over LangChain",
        "  Transparent, debuggable, minimal dependencies",
        "",
        "DuckDB over SQLite / Postgres",
        "  In-process analytical SQL, zero config, fast",
        "",
        "ChromaDB for semantic retrieval",
        "  Lightweight, embedded, cosine similarity",
        "",
        "GPT-5.4 mini with reasoning_effort=high",
        "  Best reasoning quality for drilling domain",
        "  Graceful fallback for older models",
    ], size=14)

    _add_bullets(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.5), [
        "Rule-based + LLM hybrid for phase detection",
        "  Rules: hole sizes + activity codes (deterministic)",
        "  LLM: validates against DDR narratives",
        "",
        "Domain-specific tooling over generic agents",
        "  Each tool purpose-built for a question category",
        "  Tools return structured evidence, not raw data",
        "",
        "Quality assurance",
        "  95 automated tests, 130 stress-test questions",
        "  Output validation checks all 6 answer sections",
        "  Data quality filtering at ingestion + query time",
    ], size=14)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    _accent_bar(slide, y=Inches(3.0), height=Inches(0.04), color=GOLD)

    _add_text(slide, Inches(1.2), Inches(0.6), Inches(11), Inches(0.8),
              "Why This Solution Wins", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_bullets(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2), [
        "12 DuckDB tables from 4 data sources (DDR + WITSML + production + geology)",
        "12 purpose-built agent tools covering all 6 question categories + cross-well benchmarks",
        "~38,000 searchable documents in ChromaDB vector store (DDR text + WITSML messages)",
    ], size=16, color=RGBColor(0xCC, 0xDD, 0xEE), indent_color=RGBColor(0x99, 0xBB, 0xDD))

    _add_bullets(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(3.5), [
        "Evidence-first: every answer cites BOTH data AND daily report quotes",
        "Transparent: --trace shows complete reasoning chain for every answer",
        "Domain-appropriate: hole-size phase detection, NPT sub-classification, mud correlation",
        "Deep data integration: actual measured ROP/WOB/RPM from WITSML mudlog",
        "Cross-well analysis: field-wide benchmarks, difficulty indices, gas response ranking",
        "Robust: 95 tests, 130 stress questions, data quality filtering, API retry logic",
        "Reproducible: pip install, set API key, ingest (~3 min), ask questions",
    ], size=15, color=WHITE, indent_color=RGBColor(0xBB, 0xCC, 0xDD))

    _add_text(slide, Inches(1.2), Inches(6.5), Inches(11), Inches(0.6),
              "Clean, minimal, readable code  |  No LangChain, no framework bloat",
              size=16, color=GOLD, align=PP_ALIGN.CENTER)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)           # 1
    slide_challenge(prs)       # 2
    slide_architecture(prs)    # 3
    slide_data(prs)            # 4
    slide_evidence_design(prs) # 5
    slide_example(prs)         # 6
    slide_field_benchmarks(prs)# 7
    slide_trace(prs)           # 8
    slide_design_decisions(prs)# 9
    slide_summary(prs)         # 10

    output_path = "presentation/slides.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path} ({len(prs.slides)} slides)")
    return output_path


if __name__ == "__main__":
    create_presentation()
